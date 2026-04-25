"""
build_walkthrough.py
Generates Documentation/HoneyDo_Walkthrough.pptx from the screenshots tree.
"""

import os
import re
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
BASE        = r'C:\Users\davey\source\repos\HoneyDo\Documentation\screenshots'
OUTPUT      = r'C:\Users\davey\source\repos\HoneyDo\Documentation\HoneyDo_Walkthrough.pptx'

# ---------------------------------------------------------------------------
# Brand colours
# ---------------------------------------------------------------------------
C_AMBER     = RGBColor(245, 158,  11)   # honey amber
C_AMBER_LT  = RGBColor(254, 243, 199)   # pale amber (descriptors on section slides)
C_DARK      = RGBColor( 17,  24,  39)   # near-black
C_DARK2     = RGBColor( 31,  41,  55)   # slightly lighter dark (header bar)
C_WHITE     = RGBColor(255, 255, 255)
C_SLIDE_BG  = RGBColor(248, 249, 250)   # off-white slide background
C_GRAY      = RGBColor(107, 114, 128)

# ---------------------------------------------------------------------------
# Slide geometry (16:9 widescreen)
# ---------------------------------------------------------------------------
SLIDE_W  = Inches(13.333)
SLIDE_H  = Inches(7.5)
HEADER_H = Inches(0.44)
IMG_PAD  = Inches(0.10)

# ---------------------------------------------------------------------------
# Section metadata  (folder, title, description)
# ---------------------------------------------------------------------------
SECTIONS = [
    (
        "1 - login",
        "Login & Registration",
        "Sign in, create an account, and see validation in action.",
    ),
    (
        "2 - list view",
        "List View",
        "Create, search, filter, and navigate your to-do lists.",
    ),
    (
        "3 - task view",
        "Task View",
        "Add, edit, star, sort, assign, and resolve tasks within a list — then close it.",
    ),
    (
        "4 - profile",
        "Profile",
        "Edit your profile, manage personal tags, upload an avatar, and set your colour scheme.",
    ),
    (
        "5 - friends",
        "Friends",
        "Send friend requests, accept or decline invitations, and manage your connections.",
    ),
]

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

# Hand-crafted overrides for action names that don't parse cleanly from camelCase
ACTION_OVERRIDES = {
    'failRegisterUserExists': 'Failed Registration: User Exists',
    'backtoLists':            'Back to Lists',
    'uploadAvatarUrl':        'Upload Avatar URL',
    'changePasswordGood':     'Change Password: Success',
    'changePasswordMismatch': 'Change Password: Mismatch',
    'addTaskDateAndTag':      'Add Task: Date and Tag',
    'addTaskNoDateNoTags':    'Add Task: No Date, No Tags',
    'addFriendByEmailPending':'Add Friend by Email: Pending',
    'friendListLanding':      'Friends — Landing',
}


def camel_to_words(s: str) -> str:
    """'failedLogin' → 'Failed Login'"""
    if not s:
        return s
    spaced = re.sub(r'([a-z])([A-Z])', r'\1 \2', s)
    return ' '.join(w.capitalize() for w in spaced.split())


def parse_slide_title(filename: str) -> str:
    """
    '11_starTask2_sorted.png' → 'Star Task  ·  Step 2  ·  Sorted'
    '1_landing_newUser.png'  → 'Landing  ·  New User'
    '10_openTaskList.png'    → 'Open Task List'
    """
    name = os.path.splitext(filename)[0]          # strip .png
    parts = name.split('_', 1)                    # ['11', 'starTask2_sorted']
    if len(parts) < 2:
        return camel_to_words(name)

    rest = parts[1]                               # 'starTask2_sorted'
    sub  = rest.split('_')                        # ['starTask2', 'sorted']
    main = sub[0]                                 # 'starTask2'
    descriptors = sub[1:]                         # ['sorted']

    m = re.match(r'^(.*?)(\d+)$', main)
    if m:
        action_str = m.group(1)
        step       = m.group(2)
    else:
        action_str = main
        step       = None

    # Check full action+step key first, then action alone
    full_key = action_str + (step or '')
    if full_key in ACTION_OVERRIDES:
        title = ACTION_OVERRIDES[full_key]
        step  = None   # override already includes step semantics if needed
    elif action_str in ACTION_OVERRIDES:
        title = ACTION_OVERRIDES[action_str]
    else:
        title = camel_to_words(action_str)
        if step:
            title += f'  ·  Step {step}'

    for d in descriptors:
        title += f'  ·  {camel_to_words(d)}'
    return title


def sort_key(filename: str):
    """Sort by leading integer, then alphabetically."""
    m = re.match(r'^(\d+)', filename)
    return (int(m.group(1)) if m else 999, filename)


def get_files(folder_path: str):
    return sorted(
        [f for f in os.listdir(folder_path) if f.lower().endswith('.png')],
        key=sort_key,
    )


def fit_emu(img_path: str, max_w: int, max_h: int):
    """Return (w, h) in EMU to fit image inside box, preserving aspect ratio."""
    with Image.open(img_path) as img:
        iw, ih = img.size
    ar = iw / ih
    if ar > max_w / max_h:
        w = max_w
        h = int(w / ar)
    else:
        h = max_h
        w = int(h * ar)
    return w, h


def set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color: RGBColor, line=False):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # 1 = RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line:
        shape.line.color.rgb = fill_color
    else:
        shape.line.fill.background()
    return shape


def add_label(slide, text, left, top, width, height,
              size=12, bold=False, color=C_WHITE, align=PP_ALIGN.LEFT,
              wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text       = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return txb


# ---------------------------------------------------------------------------
# Build presentation
# ---------------------------------------------------------------------------

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]  # completely blank

# ── Title slide ──────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
set_bg(slide, C_DARK)

# Amber accent bars
add_rect(slide, 0, 0, SLIDE_W, Inches(0.07), C_AMBER)
add_rect(slide, 0, SLIDE_H - Inches(0.07), SLIDE_W, Inches(0.07), C_AMBER)

# Vertical amber stripe (decorative)
add_rect(slide, Inches(0.55), Inches(1.6), Inches(0.06), Inches(4.3), C_AMBER)

# App name
add_label(slide, "HoneyDo",
    Inches(0.9), Inches(1.9), Inches(12), Inches(2.2),
    size=80, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)

# Tag line (amber)
add_label(slide, "Product Walkthrough",
    Inches(0.9), Inches(3.9), Inches(12), Inches(1.0),
    size=30, bold=False, color=C_AMBER, align=PP_ALIGN.LEFT)

# Description
add_label(slide,
    "A collaborative to-do list application  ·  "
    "Login  ·  Lists  ·  Tasks  ·  Profile  ·  Friends",
    Inches(0.9), Inches(5.0), Inches(12), Inches(0.7),
    size=14, color=C_GRAY, align=PP_ALIGN.LEFT)

# ── Sections ─────────────────────────────────────────────────────────────────
for section_folder, section_title, section_desc in SECTIONS:
    folder_path = os.path.join(BASE, section_folder)
    files       = get_files(folder_path)
    section_num = section_folder.split(' - ')[0]   # '1', '2', …

    # ── Section separator slide ───────────────────────────────────────────
    slide = prs.slides.add_slide(blank)
    set_bg(slide, C_DARK)

    # Amber left sidebar
    add_rect(slide, 0, 0, Inches(0.55), SLIDE_H, C_AMBER)

    # Large muted section number (background decoration)
    add_label(slide, section_num,
        Inches(0.7), Inches(-0.4), Inches(5), Inches(5),
        size=200, bold=True, color=RGBColor(38, 50, 68),  # very dark, barely visible
        align=PP_ALIGN.LEFT)

    # Section title
    add_label(slide, section_title,
        Inches(0.9), Inches(2.2), Inches(11.8), Inches(1.6),
        size=52, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)

    # Description
    add_label(slide, section_desc,
        Inches(0.9), Inches(3.9), Inches(11), Inches(0.8),
        size=20, color=C_AMBER_LT, align=PP_ALIGN.LEFT)

    # Count badge
    count_text = f"{len(files)} screenshot{'s' if len(files) != 1 else ''}"
    add_label(slide, count_text,
        Inches(0.9), Inches(5.0), Inches(4), Inches(0.45),
        size=13, color=C_GRAY, align=PP_ALIGN.LEFT)

    # ── Screenshot slides ─────────────────────────────────────────────────
    for filename in files:
        img_path    = os.path.join(folder_path, filename)
        slide_title = parse_slide_title(filename)

        slide = prs.slides.add_slide(blank)
        set_bg(slide, C_SLIDE_BG)

        # Header bar
        add_rect(slide, 0, 0, SLIDE_W, HEADER_H, C_DARK2)

        # Section pill label (amber text, left)
        add_label(slide,
            f"  {section_title.upper()}",
            Inches(0), Inches(0.03), Inches(4.5), HEADER_H - Inches(0.06),
            size=8, bold=True, color=C_AMBER, align=PP_ALIGN.LEFT, wrap=False)

        # Slide title (white, centred in header)
        add_label(slide, slide_title,
            Inches(3.5), Inches(0.04), Inches(6.3), HEADER_H - Inches(0.08),
            size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER, wrap=False)

        # Section progress indicator (right side of header — "N of M")
        slide_idx   = files.index(filename) + 1
        total_slides = len(files)
        add_label(slide,
            f"{slide_idx} / {total_slides}  ",
            Inches(9.5), Inches(0.05), Inches(3.7), HEADER_H - Inches(0.1),
            size=9, color=C_GRAY, align=PP_ALIGN.RIGHT, wrap=False)

        # Image — fit into remaining area
        avail_w = SLIDE_W  - 2 * IMG_PAD
        avail_h = SLIDE_H  - HEADER_H - 2 * IMG_PAD

        w, h    = fit_emu(img_path, avail_w, avail_h)

        # Centre horizontally and vertically within the available zone
        left = (SLIDE_W - w) // 2
        top  = HEADER_H + IMG_PAD + (avail_h - h) // 2

        slide.shapes.add_picture(img_path, left, top, w, h)

# ---------------------------------------------------------------------------
# Save
# ---------------------------------------------------------------------------
prs.save(OUTPUT)
print(f"Saved: {OUTPUT}")
print(f"Slides: {len(prs.slides)}")
